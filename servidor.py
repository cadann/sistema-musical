"""
servidor.py — Musica Facil v5.3
Rotas:
  GET  /               — serve index.html (abra http://localhost:5000)
  POST /api/converter  — converte PPTX/PPT via conversor_pptx.py
  POST /api/pdf        — extrai musicas de PDF via Claude Opus 4.6
  GET  /api/status     — status do servidor
"""

from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from flask_socketio import SocketIO, emit, join_room
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os, re, unicodedata, traceback, time, tempfile, sys, base64, json
from pathlib import Path

app = Flask(__name__)
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'musicafacil2025')
CORS(app, resources={r"/*": {"origins": "*"}})
socketio = SocketIO(app, cors_allowed_origins="*")

# ── BANCO DE DADOS (PostgreSQL na nuvem, SQLite local) ────────────────────────
DATABASE_URL = os.environ.get('DATABASE_URL')

if DATABASE_URL:
    # PostgreSQL no Render
    import psycopg2
    import psycopg2.extras
    # Render usa "postgres://" mas psycopg2 precisa de "postgresql://"
    if DATABASE_URL.startswith('postgres://'):
        DATABASE_URL = DATABASE_URL.replace('postgres://', 'postgresql://', 1)

    def get_db():
        conn = psycopg2.connect(DATABASE_URL)
        return conn

    def init_db():
        conn = get_db()
        c = conn.cursor()
        c.execute("""
            CREATE TABLE IF NOT EXISTS musicas (
                id BIGINT PRIMARY KEY,
                title TEXT NOT NULL,
                artist TEXT,
                key TEXT,
                tags TEXT,
                cifra TEXT,
                strophes TEXT,
                obs TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        c.execute("""
            CREATE TABLE IF NOT EXISTS playlists (
                id BIGINT PRIMARY KEY,
                name TEXT NOT NULL,
                desc TEXT,
                songs TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        c.execute("""
            CREATE TABLE IF NOT EXISTS membros (
                id BIGINT PRIMARY KEY,
                tipo TEXT,
                nome TEXT NOT NULL,
                funcao TEXT,
                contato TEXT,
                obs TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        c.execute("""
            CREATE TABLE IF NOT EXISTS eventos (
                id BIGINT PRIMARY KEY,
                nome TEXT NOT NULL,
                local TEXT,
                data TEXT,
                hora TEXT,
                duracao TEXT,
                playlist_id BIGINT,
                musico_id BIGINT,
                co_musicos_ids TEXT,
                operador_id BIGINT,
                codigo TEXT UNIQUE,
                status TEXT DEFAULT 'upcoming',
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        conn.commit()
        conn.close()
        print('  Banco PostgreSQL inicializado.')

    def executar(sql, params=()):
        sql = sql.replace('?', '%s')
        sql = sql.replace('INSERT OR REPLACE', 'INSERT')
        sql = sql.replace('last_insert_rowid()', 'lastval()')
        return sql, params

    def rows_to_list(cursor):
        cols = [d[0] for d in cursor.description]
        return [dict(zip(cols, row)) for row in cursor.fetchall()]

    def row_to_dict(cursor):
        cols = [d[0] for d in cursor.description]
        row = cursor.fetchone()
        return dict(zip(cols, row)) if row else None

else:
    # SQLite local
    import sqlite3
    DB_PATH = Path(__file__).parent.resolve() / 'musicafacil.db'

    def get_db():
        conn = sqlite3.connect(str(DB_PATH))
        conn.row_factory = sqlite3.Row
        return conn

    def init_db():
        conn = get_db()
        conn.executescript("""
            CREATE TABLE IF NOT EXISTS musicas (
                id INTEGER PRIMARY KEY, title TEXT NOT NULL, artist TEXT,
                key TEXT, tags TEXT, cifra TEXT, strophes TEXT, obs TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            CREATE TABLE IF NOT EXISTS playlists (
                id INTEGER PRIMARY KEY, name TEXT NOT NULL, desc TEXT, songs TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            CREATE TABLE IF NOT EXISTS membros (
                id INTEGER PRIMARY KEY, tipo TEXT, nome TEXT NOT NULL,
                funcao TEXT, contato TEXT, obs TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            CREATE TABLE IF NOT EXISTS eventos (
                id INTEGER PRIMARY KEY, nome TEXT NOT NULL, local TEXT,
                data TEXT, hora TEXT, duracao TEXT, playlist_id INTEGER,
                musico_id INTEGER, co_musicos_ids TEXT, operador_id INTEGER,
                codigo TEXT UNIQUE, status TEXT DEFAULT 'upcoming',
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
        """)
        conn.commit()
        conn.close()
        print('  Banco SQLite inicializado:', DB_PATH)

    def executar(sql, params=()):
        return sql, params

    def rows_to_list(rows):
        return [dict(r) for r in rows]

    def row_to_dict(row):
        return dict(row) if row else None

init_db()

# ── ESTADO COMPARTILHADO ──────────────────────────────────────────────────────
estado = {
    'songIdx': 0,
    'stropheIdx': 0,
    'musicoConectado': False,
}

PASTA_TXT = Path(r"C:\Musicas_pptx")


# ── ROTAS PRINCIPAIS ───────────────────────────────────────────────────────────

def _serve_index():
    for base in (Path(__file__).parent.resolve(), Path.cwd()):
        if (base / 'index.html').exists():
            return send_from_directory(str(base), 'index.html')
    return "<h2>index.html nao encontrado.</h2>", 404

@app.route('/')
def index():
    return _serve_index()

@app.route('/evento/<codigo>')
def evento_comusico(codigo):
    """Co-Músico acessa via QR Code — carrega o app com o código do evento."""
    # Serve o mesmo index.html; o JS detecta o path e entra como Co-Músico
    return _serve_index()


# ── UTILITÁRIOS ───────────────────────────────────────────────────────────────

def limpar_texto(texto: str) -> str:
    if not texto:
        return texto
    substituicoes = [
        ('Ã£','a'),('Ã§','c'),('Ã©','e'),('Ã­','i'),('Ã³','o'),('Ãº','u'),
        ('Ã¡','a'),('Ã ','a'),('Ã¢','a'),('Ãª','e'),('Ã´','o'),('Ã»','u'),
        ('Ã•','O'),('Ã‡','C'),('Ã‰','E'),('Ã"','O'),('Ã€','A'),('Ã‚','A'),
        ('Ãœ','U'),('A£','a'),('A§','c'),('A©','e'),('A­','i'),('A³','o'),
        ('Aº','u'),('A¡','a'),('A¼','u'),('A‡','C'),('Aª','a'),('A°','o'),
    ]
    for errado, correto in substituicoes:
        texto = texto.replace(errado, correto)
    return ''.join(
        c for c in unicodedata.normalize('NFD', texto)
        if unicodedata.category(c) != 'Mn'
    )


_CHORD_PAT = re.compile(
    r'\b[A-G][#b]?'
    r'(?:maj\d*|min|m(?:aj)?\d*|dim\d*|aug|sus[24]?|add\d+|\d+)*'
    r'(?:/[A-G][#b]?)?\b'
)
_LABEL_PREFIX = re.compile(r'^(?:Introd?\w*|Intro|Tom|Ton|Coda|Refrao)\s*[:\.]?\s*', flags=re.IGNORECASE)
_LEFTOVER = re.compile(r'[\s()\[\]x0-9+\-/#_.,;:!?|\'\"]+')

def is_chord_line(linha: str) -> bool:
    s = linha.strip()
    if not s: return False
    s = _LABEL_PREFIX.sub('', s).strip()
    s2 = _CHORD_PAT.sub('', s)
    s2 = _LEFTOVER.sub('', s2)
    return len(s2) <= 2


# ── ROTAS REST — MÚSICAS ─────────────────────────────────────────────────────

@app.route('/api/musicas', methods=['GET'])
def get_musicas():
    conn = get_db()
    c = conn.cursor()
    c.execute('SELECT * FROM musicas ORDER BY id DESC')
    musicas = []
    for r in rows_to_list(c):
        r['tags'] = json.loads(r['tags'] or '[]')
        r['strophes'] = json.loads(r['strophes'] or '[]')
        musicas.append(r)
    conn.close()
    return jsonify(musicas)

@app.route('/api/musicas', methods=['POST'])
def save_musica():
    data = request.json
    conn = get_db()
    c = conn.cursor()
    sql, params = executar(
        'INSERT OR REPLACE INTO musicas (id,title,artist,key,tags,cifra,strophes,obs) VALUES (?,?,?,?,?,?,?,?)',
        (data.get('id'), data['title'], data.get('artist',''), data.get('key','?'),
         json.dumps(data.get('tags',[])), data.get('cifra',''),
         json.dumps(data.get('strophes',[])), data.get('obs',''))
    )
    c.execute(sql, params)
    conn.commit()
    conn.close()
    return jsonify({'ok': True, 'id': data.get('id')})

@app.route('/api/musicas/bulk', methods=['POST'])
def save_musicas_bulk():
    musicas = request.json
    conn = get_db()
    c = conn.cursor()
    for data in musicas:
        sql, params = executar(
            'INSERT OR REPLACE INTO musicas (id,title,artist,key,tags,cifra,strophes,obs) VALUES (?,?,?,?,?,?,?,?)',
            (data.get('id'), data['title'], data.get('artist',''), data.get('key','?'),
             json.dumps(data.get('tags',[])), data.get('cifra',''),
             json.dumps(data.get('strophes',[])), data.get('obs',''))
        )
        c.execute(sql, params)
    conn.commit()
    conn.close()
    return jsonify({'ok': True, 'total': len(musicas)})

@app.route('/api/musicas/<int:mid>', methods=['DELETE'])
def delete_musica(mid):
    conn = get_db()
    c = conn.cursor()
    sql, params = executar('DELETE FROM musicas WHERE id=?', (mid,))
    c.execute(sql, params)
    conn.commit()
    conn.close()
    return jsonify({'ok': True})

# ── ROTAS REST — PLAYLISTS ────────────────────────────────────────────────────

@app.route('/api/playlists', methods=['GET'])
def get_playlists():
    conn = get_db()
    c = conn.cursor()
    c.execute('SELECT * FROM playlists ORDER BY id DESC')
    result = []
    for r in rows_to_list(c):
        r['songs'] = json.loads(r['songs'] or '[]')
        result.append(r)
    conn.close()
    return jsonify(result)

@app.route('/api/playlists', methods=['POST'])
def save_playlist():
    data = request.json
    conn = get_db()
    c = conn.cursor()
    sql, params = executar(
        'INSERT OR REPLACE INTO playlists (id,name,desc,songs) VALUES (?,?,?,?)',
        (data.get('id'), data['name'], data.get('desc',''), json.dumps(data.get('songs',[])))
    )
    c.execute(sql, params)
    conn.commit()
    conn.close()
    return jsonify({'ok': True, 'id': data.get('id')})

@app.route('/api/playlists/<int:pid>', methods=['DELETE'])
def delete_playlist(pid):
    conn = get_db()
    c = conn.cursor()
    sql, params = executar('DELETE FROM playlists WHERE id=?', (pid,))
    c.execute(sql, params)
    conn.commit()
    conn.close()
    return jsonify({'ok': True})

# ── ROTAS REST — MEMBROS ──────────────────────────────────────────────────────

@app.route('/api/membros', methods=['GET'])
def get_membros():
    conn = get_db()
    c = conn.cursor()
    c.execute('SELECT * FROM membros ORDER BY nome')
    result = rows_to_list(c)
    conn.close()
    return jsonify(result)

@app.route('/api/membros', methods=['POST'])
def save_membro():
    data = request.json
    conn = get_db()
    c = conn.cursor()
    sql, params = executar(
        'INSERT OR REPLACE INTO membros (id,tipo,nome,funcao,contato,obs) VALUES (?,?,?,?,?,?)',
        (data.get('id'), data.get('tipo','musico'), data['nome'],
         data.get('funcao',''), data.get('contato',''), data.get('obs',''))
    )
    c.execute(sql, params)
    conn.commit()
    conn.close()
    return jsonify({'ok': True, 'id': data.get('id')})

@app.route('/api/membros/<int:mid>', methods=['DELETE'])
def delete_membro(mid):
    conn = get_db()
    c = conn.cursor()
    sql, params = executar('DELETE FROM membros WHERE id=?', (mid,))
    c.execute(sql, params)
    conn.commit()
    conn.close()
    return jsonify({'ok': True})

# ── ROTAS REST — EVENTOS ──────────────────────────────────────────────────────

@app.route('/api/eventos', methods=['GET'])
def get_eventos():
    conn = get_db()
    c = conn.cursor()
    c.execute('SELECT * FROM eventos ORDER BY data DESC, hora DESC')
    result = []
    for r in rows_to_list(c):
        r['coMusicosIds'] = json.loads(r.pop('co_musicos_ids') or '[]')
        r['playlistId']   = r.pop('playlist_id')
        r['musicoId']     = r.pop('musico_id')
        r['operadorId']   = r.pop('operador_id')
        result.append(r)
    conn.close()
    return jsonify(result)

@app.route('/api/eventos', methods=['POST'])
def save_evento():
    data = request.json
    conn = get_db()
    c = conn.cursor()
    sql, params = executar(
        'INSERT OR REPLACE INTO eventos (id,nome,local,data,hora,duracao,playlist_id,musico_id,co_musicos_ids,operador_id,codigo,status) VALUES (?,?,?,?,?,?,?,?,?,?,?,?)',
        (data.get('id'), data['nome'], data.get('local',''),
         data.get('data',''), data.get('hora',''), data.get('duracao',''),
         data.get('playlistId'), data.get('musicoId'),
         json.dumps(data.get('coMusicosIds',[])),
         data.get('operadorId') or None,
         data.get('codigo'), data.get('status','upcoming'))
    )
    c.execute(sql, params)
    conn.commit()
    conn.close()
    return jsonify({'ok': True, 'id': data.get('id')})

@app.route('/api/eventos/<int:eid>', methods=['DELETE'])
def delete_evento(eid):
    conn = get_db()
    c = conn.cursor()
    sql, params = executar('DELETE FROM eventos WHERE id=?', (eid,))
    c.execute(sql, params)
    conn.commit()
    conn.close()
    return jsonify({'ok': True})

@app.route('/api/eventos/codigo/<codigo>', methods=['GET'])
def get_evento_por_codigo(codigo):
    conn = get_db()
    c = conn.cursor()
    sql, params = executar('SELECT * FROM eventos WHERE codigo=?', (codigo,))
    c.execute(sql, params)
    result = rows_to_list(c)
    conn.close()
    if not result:
        return jsonify({'erro': 'Evento nao encontrado'}), 404
    e = result[0]
    e['coMusicosIds'] = json.loads(e.pop('co_musicos_ids') or '[]')
    e['playlistId']   = e.pop('playlist_id')
    e['musicoId']     = e.pop('musico_id')
    e['operadorId']   = e.pop('operador_id')
    return jsonify(e)


# ── PPTX → SONG ──────────────────────────────────────────────────────────────

def extrair_linhas_slide(slide) -> list:
    linhas = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            pass
        for para in shape.text_frame.paragraphs:
            linha = ''.join(run.text for run in para.runs)
            linha = limpar_texto(linha)
            if linha.strip():
                linhas.append(linha)
    return linhas


def montar_bloco(linhas: list) -> str:
    if not linhas: return ''
    resultado = []
    encontrou_cifra = False
    ultimo_era_letra = False
    for linha in linhas:
        eh_cifra = is_chord_line(linha)
        if eh_cifra:
            encontrou_cifra = True
            resultado.append(linha)
            ultimo_era_letra = False
        else:
            if encontrou_cifra and ultimo_era_letra and resultado:
                resultado[-1] = resultado[-1] + ' ' + linha.strip()
            else:
                resultado.append(linha)
            ultimo_era_letra = True
    return '\n'.join(resultado)


def converter_ppt_para_pptx(caminho_ppt: Path, tmpdir: str) -> Path:
    soffice = r"C:\Program Files\LibreOffice\program\soffice.exe" if sys.platform == "win32" else "libreoffice"
    import subprocess
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pptx", "--outdir", tmpdir, str(caminho_ppt)],
        capture_output=True,
    )
    resultado = Path(tmpdir) / (caminho_ppt.stem + ".pptx")
    if not resultado.exists():
        raise RuntimeError(f"LibreOffice nao converteu {caminho_ppt.name}")
    return resultado


def pptx_para_song(nome_arquivo: str, prs: Presentation) -> dict:
    titulo = limpar_texto(Path(nome_arquivo).stem)
    titulo = re.sub(r'[_\-]+', ' ', titulo).strip()
    artist = '-'
    blocos = []

    for i, slide in enumerate(prs.slides):
        linhas = extrair_linhas_slide(slide)
        if not linhas: continue
        if i == 0:
            nao_vazias = [l for l in linhas if l.strip()]
            if nao_vazias and not re.match(r'^[A-G]', nao_vazias[0]):
                titulo = nao_vazias[0].strip()
            for linha in nao_vazias:
                m = re.match(r'^\((.+)\)$', linha.strip())
                if m: artist = m.group(1).strip(); break
            uteis = [l for l in nao_vazias if l.strip() != titulo and not re.match(r'^\(.+\)$', l.strip())]
            if not uteis: continue
        bloco = montar_bloco(list(linhas))
        if bloco.strip(): blocos.append(bloco)

    if not blocos: blocos = ['(sem conteudo)']

    cifra = '\n\n'.join(blocos)

    try:
        PASTA_TXT.mkdir(parents=True, exist_ok=True)
        nome_txt = limpar_texto(Path(nome_arquivo).stem).replace(' ', '_') + '.txt'
        (PASTA_TXT / nome_txt).write_text(titulo + '\n\n' + cifra, encoding='utf-8')
    except Exception as e:
        print(f'  Aviso TXT: {e}')

    def extrair_letra(bloco):
        return '\n'.join(l.strip() for l in bloco.split('\n') if l.strip() and not is_chord_line(l))

    strophes = [extrair_letra(b) for b in blocos]
    strophes = [s for s in strophes if s.strip()] or ['(sem conteudo)']

    return {
        'id': int(time.time() * 1000) + len(strophes),
        'title': titulo, 'artist': artist, 'key': '?', 'tags': [],
        'cifra': cifra, 'strophes': strophes,
    }


# ── ROTAS ─────────────────────────────────────────────────────────────────────

@app.route('/api/converter', methods=['POST'])
def converter():
    arquivos = request.files.getlist('arquivos')
    if not arquivos:
        return jsonify({'erro': 'Nenhum arquivo recebido.'}), 400
    importadas, erros = [], []
    for arquivo in arquivos:
        nome = arquivo.filename or 'musica.pptx'
        try:
            if nome.lower().endswith('.ppt'):
                with tempfile.TemporaryDirectory() as tmpdir:
                    caminho_ppt = Path(tmpdir) / nome
                    arquivo.save(str(caminho_ppt))
                    caminho_pptx = converter_ppt_para_pptx(caminho_ppt, tmpdir)
                    prs = Presentation(str(caminho_pptx))
            else:
                prs = Presentation(arquivo.stream)
            song = pptx_para_song(nome, prs)
            importadas.append(song)
            print(f'  OK: {song["title"]} — {len(song["strophes"])} estrofes')
        except Exception as e:
            erros.append({'arquivo': nome, 'erro': str(e)})
            traceback.print_exc()
    return jsonify({'importadas': importadas, 'erros': erros})


# ── REORGANIZAR COLUNAS ───────────────────────────────────────────────────────

def reorganizar_colunas(texto):
    """
    Detecta paginas com layout de duas colunas e reorganiza:
    coluna esquerda completa primeiro, depois coluna direita.
    """
    sep_pagina = '\f'
    paginas_saida = []

    for pagina in texto.split(sep_pagina):
        linhas = pagina.split('\n')
        if not any(l.strip() for l in linhas):
            continue

        largura = max((len(l) for l in linhas), default=0)
        if largura < 40:
            paginas_saida.append(pagina)
            continue

        meio = largura // 2
        linhas_com_conteudo = [l for l in linhas if l.strip()]
        if not linhas_com_conteudo:
            paginas_saida.append(pagina)
            continue

        duas_colunas = sum(
            1 for l in linhas_com_conteudo
            if len(l) > meio and l[:meio].strip() and l[meio:].strip()
        )

        if duas_colunas / len(linhas_com_conteudo) < 0.4:
            paginas_saida.append(pagina)
            continue

        col_esq = []
        col_dir = []

        for linha in linhas:
            esq = linha[:meio].rstrip()
            dir_ = linha[meio:].rstrip()
            if esq.strip() and dir_.strip():
                col_esq.append(esq)
                col_dir.append(dir_.lstrip())
            elif esq.strip():
                col_esq.append(esq)
            elif dir_.strip():
                col_dir.append(dir_.lstrip())
            else:
                col_esq.append('')
                col_dir.append('')

        while col_esq and not col_esq[-1].strip():
            col_esq.pop()
        while col_dir and not col_dir[-1].strip():
            col_dir.pop()

        reorganizado = '\n'.join(col_esq) + '\n\n' + '\n'.join(col_dir)
        paginas_saida.append(reorganizado)
        print('  Pagina com 2 colunas detectada e reorganizada.')

    return sep_pagina.join(paginas_saida)


@app.route('/api/pdf', methods=['POST'])
def importar_pdf():
    """
    Recebe:
      - arquivo PDF (multipart 'arquivo')
      - claude_key (form field)
    Chama Claude Opus 4.6 e retorna lista de músicas.
    """
    arquivo = request.files.get('arquivo')
    claude_key = request.form.get('claude_key', '').strip()

    if not arquivo:
        return jsonify({'erro': 'Nenhum arquivo recebido.'}), 400
    if not claude_key:
        return jsonify({'erro': 'Claude API Key não fornecida.'}), 400

    try:
        import requests as req
        import subprocess

        # ── Extrai texto do PDF localmente com pdftotext ──────────────────────
        # Muito mais rápido que enviar o PDF inteiro em base64 para a API
        pdf_bytes = arquivo.read()

        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(pdf_bytes)
            tmp_path = tmp.name

        try:
            result = subprocess.run(
                ['pdftotext', '-layout', '-enc', 'UTF-8', tmp_path, '-'],
                capture_output=True, timeout=30
            )
            texto_pdf = result.stdout.decode('utf-8', errors='replace')
            texto_pdf = reorganizar_colunas(texto_pdf)
        except (subprocess.TimeoutExpired, FileNotFoundError):
            # pdftotext não disponível: fallback para base64
            texto_pdf = None
        finally:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass

        prompt = """Você é um especialista em cifras musicais brasileiras.
Analise este hinário musical e extraia TODAS as músicas.

REGRAS:
1. O padrão de início de nova música é: NOME DO ARTISTA (em maiúsculas) seguido do título na linha seguinte.
2. Se a primeira linha de uma página contém acordes musicais (ex: F C/E, G Em, D7M), é CONTINUAÇÃO da música anterior.
3. Ignore a palavra "ÍNDICE" que aparece nos cabeçalhos.
4. Preserve os acentos do português (ã, ç, é, ê, ô, etc.).
5. Preserve os acordes exatamente como estão (ex: F#m7(9), D7M, Bm7(b5)).
6. ESPAÇAMENTO DOS ACORDES — regra crítica:
   - As linhas de acorde têm espaços INICIAIS que alinham cada acorde exatamente acima da sílaba correspondente.
   - Exemplo correto:
         F#m              D
     Me perguntam de onde vem
   - Neste exemplo, "F#m" está recuado com espaços para ficar acima de "perguntam", e "D" está acima de "vem".
   - Você DEVE copiar esses espaços iniciais e intermediários EXATAMENTE como estão no original.
   - NUNCA remova, reduza ou normalize os espaços de uma linha de acorde.
   - NUNCA mova um acorde para o início da linha se ele não estava lá no original.

Retorne APENAS um array JSON válido, sem explicações, sem markdown.
Cada item deve ter:
- "title": título da música
- "artist": nome do artista
- "cifra": texto completo com acordes e letra (preservando todos os espaços e acentos exatamente como no original)
- "strophes": array de strings com cada estrofe SEM acordes (apenas a letra)"""

        # Monta a mensagem: texto puro (rápido) ou base64 (fallback)
        if texto_pdf and texto_pdf.strip():
            print(f'  PDF → texto local ({len(texto_pdf)} chars), enviando só texto para Claude...')
            content_msg = [{'type': 'text', 'text': f"{prompt}\n\n---\n{texto_pdf}"}]
        else:
            print('  pdftotext indisponível, enviando PDF em base64...')
            pdf_b64 = base64.b64encode(pdf_bytes).decode('utf-8')
            content_msg = [
                {'type': 'document', 'source': {'type': 'base64', 'media_type': 'application/pdf', 'data': pdf_b64}},
                {'type': 'text', 'text': prompt}
            ]

        response = req.post(
            'https://api.anthropic.com/v1/messages',
            headers={
                'x-api-key': claude_key,
                'anthropic-version': '2023-06-01',
                'Content-Type': 'application/json',
            },
            json={
                'model': 'claude-sonnet-4-6',
                'max_tokens': 8000,
                'messages': [{'role': 'user', 'content': content_msg}]
            },
            timeout=300
        )

        if response.status_code != 200:
            # Extrai mensagem legível da resposta da API
            try:
                err_data = response.json()
                err_msg = (err_data.get('error', {}) or {}).get('message', '')
            except Exception:
                err_msg = response.text[:300]

            # Traduz erros comuns para mensagens amigáveis
            if 'credit balance' in err_msg.lower() or 'billing' in err_msg.lower():
                msg_amigavel = ('Saldo insuficiente na Anthropic. '
                                'Adicione créditos em console.anthropic.com/settings/billing')
            elif 'invalid_api_key' in err_msg.lower() or 'authentication' in err_msg.lower():
                msg_amigavel = 'Claude API Key inválida. Verifique a chave na barra lateral.'
            elif 'overloaded' in err_msg.lower():
                msg_amigavel = 'A API do Claude está sobrecarregada. Tente novamente em alguns segundos.'
            elif err_msg:
                msg_amigavel = f'Erro da API Claude: {err_msg}'
            else:
                msg_amigavel = f'Erro na API Claude (HTTP {response.status_code}). Verifique sua chave e créditos.'

            return jsonify({'erro': msg_amigavel}), 500

        data = response.json()
        texto = data['content'][0]['text'].strip()
        texto = re.sub(r'^```json\s*', '', texto, flags=re.IGNORECASE)
        texto = re.sub(r'^```\s*', '', texto, flags=re.IGNORECASE)
        texto = re.sub(r'\s*```$', '', texto)
        texto = texto.strip()

        # Tenta parsear normalmente
        try:
            musicas = json.loads(texto)
        except json.JSONDecodeError as je:
            # JSON cortado: tenta recuperar músicas completas até onde veio
            print(f'  Aviso: JSON incompleto ({je}), tentando recuperar...')
            # Encontra o último objeto completo antes do corte
            ultimo_ok = texto.rfind('},')
            if ultimo_ok == -1:
                ultimo_ok = texto.rfind('}')
            if ultimo_ok > 0:
                parcial = texto[:ultimo_ok+1].strip()
                # Garante que é um array válido
                if not parcial.startswith('['):
                    parcial = '[' + parcial
                parcial += ']'
                try:
                    musicas = json.loads(parcial)
                    print(f'  Recuperadas {len(musicas)} música(s) do JSON parcial')
                except Exception:
                    raise ValueError(
                        f'Resposta da API incompleta (JSON cortado). '
                        f'O PDF pode ser muito longo. Tente dividir em partes menores.'
                    )
            else:
                raise ValueError(
                    'Resposta da API vazia ou inválida. Tente novamente.'
                )

        print(f'  PDF: {len(musicas)} música(s) extraída(s)')
        return jsonify({'musicas': musicas})

    except Exception as e:
        traceback.print_exc()
        err_str = str(e)
        if 'timed out' in err_str.lower() or 'timeout' in err_str.lower():
            msg = ('O PDF demorou demais para ser processado (timeout). '
                   'Tente um PDF menor ou com menos páginas.')
        elif 'connection' in err_str.lower() or 'network' in err_str.lower():
            msg = 'Erro de conexão com a API Claude. Verifique sua internet e tente novamente.'
        else:
            msg = err_str
        return jsonify({'erro': msg}), 500


@app.route('/api/status', methods=['GET'])
def status():
    return jsonify({'status': 'ok', 'versao': '5.4', 'estado': estado})


# ── WEBSOCKET EVENTOS ─────────────────────────────────────────────────────────

@socketio.on('connect')
def on_connect():
    print(f'  Cliente conectado: {request.sid}')
    # Envia estado atual para o cliente que acabou de conectar
    emit('estado', estado)

@socketio.on('disconnect')
def on_disconnect():
    papel = request.environ.get('HTTP_X_PAPEL', '?')
    if papel == 'musico':
        estado['musicoConectado'] = False
        print('  Músico desconectado.')

@socketio.on('registrar')
def on_registrar(data):
    papel = data.get('papel', '')
    request.environ['HTTP_X_PAPEL'] = papel
    if papel == 'musico':
        estado['musicoConectado'] = True
        print(f'  Músico conectado.')
    emit('estado', estado)

@socketio.on('navegar')
def on_navegar(data):
    papel = data.get('papel', '')
    # Hierarquia: músico sempre passa, operador só se músico ausente
    if papel == 'comusico':
        return  # Co-Músico não pode navegar
    if papel == 'operador' and estado['musicoConectado']:
        return  # Operador bloqueado se músico está conectado

    estado['songIdx']    = data.get('songIdx', estado['songIdx'])
    estado['stropheIdx'] = data.get('stropheIdx', estado['stropheIdx'])

    # Retransmite para TODOS os clientes
    emit('estado', estado, broadcast=True)
    print(f'  Navegação por {papel}: música {estado["songIdx"]}, estrofe {estado["stropheIdx"]}')


if __name__ == '__main__':
    print('=' * 55)
    print('  Musica Facil — Servidor v5.4')
    print('  PPTX: python-pptx + LibreOffice')
    print('  PDF:  Claude Sonnet')
    print('  WebSocket: flask-socketio')
    print('=' * 55)
    print('  Acesso local:  http://localhost:5000')
    print('  Acesso na rede: http://SEU-IP:5000')
    print('  (NAO abra o index.html direto como arquivo)')
    print('  Ctrl+C para encerrar')
    print('=' * 55)
    socketio.run(app, host='0.0.0.0', port=5000, debug=False, allow_unsafe_werkzeug=True)
